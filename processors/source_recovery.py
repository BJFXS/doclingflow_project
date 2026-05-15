from __future__ import annotations

"""Source-aware text recovery helpers for fidelity-sensitive conversions."""

from dataclasses import dataclass
from difflib import SequenceMatcher
from pathlib import Path
import re

from pipeline.ocr_quality import collect_ocr_page_risk_summary, page_quality_metrics, page_quality_score, retry_page_passes_absolute_checks


ELLIPSIS_RE = re.compile(r"(?:\.{3,}|…)")
MARKDOWN_PREFIX_RE = re.compile(r"^(?P<prefix>\s*(?:[-*+]\s+|#+\s+|\d+\.\s+|>\s+)?)(?P<content>.*)$")
ALNUM_RE = re.compile(r"[A-Za-z0-9]{3,}")


@dataclass(frozen=True)
class SourceRecoveryResult:
    """Recovered Markdown plus lightweight accounting for reports and logs."""

    markdown: str
    replaced_line_count: int = 0
    appended_line_count: int = 0


@dataclass(frozen=True)
class SlideTextLine:
    """One normalized text line extracted from a PowerPoint slide."""

    slide_number: int
    text: str


def recover_pptx_markdown(markdown: str, source_path: Path) -> SourceRecoveryResult:
    """Repair obvious PPTX text damage and append any still-missing slide text."""

    slide_lines = _extract_pptx_slide_lines(source_path)
    if not slide_lines:
        return SourceRecoveryResult(markdown=markdown)

    replaced_markdown, replaced_line_count = _replace_degraded_markdown_lines(markdown, slide_lines)
    missing_lines = _collect_missing_slide_lines(replaced_markdown, slide_lines)
    if not missing_lines:
        return SourceRecoveryResult(markdown=replaced_markdown, replaced_line_count=replaced_line_count)

    recovered = replaced_markdown.rstrip()
    if recovered:
        recovered += "\n\n"
    recovered += "## Slide Text Recovery\n\n"
    recovered += "> The lines below were preserved directly from the source slides because the Markdown export looked incomplete.\n\n"

    current_slide = None
    appended_line_count = 0
    for line in missing_lines:
        if line.slide_number != current_slide:
            if current_slide is not None:
                recovered += "\n"
            current_slide = line.slide_number
            recovered += f"### Slide {current_slide}\n\n"
        recovered += f"- {line.text}\n"
        appended_line_count += 1
    recovered = recovered.rstrip() + "\n"
    return SourceRecoveryResult(
        markdown=recovered,
        replaced_line_count=replaced_line_count,
        appended_line_count=appended_line_count,
    )


def append_pdf_text_layer_recovery(
    markdown: str,
    source_path: Path,
    source_pages: list[str],
    strategy: object,
) -> SourceRecoveryResult:
    """Append higher-confidence text-layer excerpts for OCR-risky PDF pages."""

    if source_path.suffix.lower() != ".pdf" or not getattr(strategy, "enable_ocr", False):
        return SourceRecoveryResult(markdown=markdown)
    risk_summary = collect_ocr_page_risk_summary(source_pages, strategy)
    risky_pages = list(risk_summary.get("risky_pages") or [])
    if not risky_pages:
        return SourceRecoveryResult(markdown=markdown)

    raw_pages = _extract_pdf_text_pages(source_path)
    if not raw_pages:
        return SourceRecoveryResult(markdown=markdown)

    recovered_pages: list[tuple[int, str]] = []
    for page_number in risky_pages:
        if page_number < 1 or page_number > len(raw_pages) or page_number > len(source_pages):
            continue
        raw_page = raw_pages[page_number - 1].strip()
        ocr_page = source_pages[page_number - 1].strip()
        if not raw_page or not retry_page_passes_absolute_checks(raw_page):
            continue
        raw_score = page_quality_score(page_quality_metrics(raw_page))
        ocr_score = page_quality_score(page_quality_metrics(ocr_page))
        if raw_score <= ocr_score + 0.12:
            continue
        recovered_pages.append((page_number, raw_page))

    if not recovered_pages:
        return SourceRecoveryResult(markdown=markdown)

    recovered = markdown.rstrip()
    if recovered:
        recovered += "\n\n"
    recovered += "### Higher-Confidence Text Layer Recovery\n\n"
    recovered += "> The pages below had risky OCR output. A stronger embedded text layer was also preserved when available.\n\n"
    for page_number, page_text in recovered_pages:
        recovered += f"#### Page {page_number}\n\n{page_text}\n\n"
    recovered = recovered.rstrip() + "\n"
    return SourceRecoveryResult(markdown=recovered, appended_line_count=len(recovered_pages))


def _extract_pptx_slide_lines(source_path: Path) -> list[SlideTextLine]:
    """Extract paragraph-level text from a PPTX while preserving slide numbers."""

    try:
        from pptx import Presentation
    except Exception:
        return []

    try:
        presentation = Presentation(str(source_path))
    except Exception:
        return []

    slide_lines: list[SlideTextLine] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            for paragraph_text in _extract_shape_text(shape):
                normalized = _normalize_visible_line(paragraph_text)
                if normalized:
                    slide_lines.append(SlideTextLine(slide_number=slide_index, text=normalized))
    return slide_lines


def _extract_shape_text(shape: object) -> list[str]:
    """Collect visible text from a slide shape, including grouped children."""

    lines: list[str] = []
    shapes = getattr(shape, "shapes", None)
    if shapes is not None:
        for child in shapes:
            lines.extend(_extract_shape_text(child))
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return lines
    for paragraph in getattr(text_frame, "paragraphs", ()):
        parts = [run.text for run in getattr(paragraph, "runs", ()) if getattr(run, "text", "").strip()]
        if not parts:
            raw_text = getattr(paragraph, "text", "")
            if raw_text.strip():
                parts = [raw_text]
        merged = "".join(parts).strip()
        if merged:
            lines.append(merged)
    return lines


def _replace_degraded_markdown_lines(markdown: str, slide_lines: list[SlideTextLine]) -> tuple[str, int]:
    """Replace obviously degraded Markdown lines with a better slide-text match."""

    replaced_count = 0
    replaced_lines: list[str] = []
    for line in markdown.splitlines():
        match = MARKDOWN_PREFIX_RE.match(line)
        if match is None:
            replaced_lines.append(line)
            continue
        prefix = match.group("prefix")
        content = match.group("content").strip()
        best = None
        if _looks_like_degraded_pptx_line(content):
            best = _best_source_line_match(content, slide_lines)
        elif _looks_like_case_softened_source_line(content):
            best = _best_source_line_case_refresh(content, slide_lines)
        if best is None:
            replaced_lines.append(line)
            continue
        replaced_lines.append(f"{prefix}{best.text}".rstrip())
        replaced_count += 1
    return "\n".join(replaced_lines), replaced_count


def _collect_missing_slide_lines(markdown: str, slide_lines: list[SlideTextLine]) -> list[SlideTextLine]:
    """Return substantive source slide lines that are still missing from Markdown."""

    normalized_markdown_lines = [_normalize_markdown_line(line) for line in markdown.splitlines()]
    missing: list[SlideTextLine] = []
    for line in slide_lines:
        if len(ALNUM_RE.findall(line.text)) < 2:
            continue
        if _markdown_already_covers_line(line.text, normalized_markdown_lines):
            continue
        missing.append(line)
    return missing


def _markdown_already_covers_line(source_line: str, normalized_markdown_lines: list[str]) -> bool:
    """Decide whether the existing Markdown already covers one source line."""

    target = _normalize_visible_line(source_line).casefold()
    if not target:
        return True
    for markdown_line in normalized_markdown_lines:
        if not markdown_line:
            continue
        if target in markdown_line or markdown_line in target:
            return True
        if SequenceMatcher(None, target, markdown_line).ratio() >= 0.92:
            return True
    return False


def _best_source_line_match(markdown_content: str, slide_lines: list[SlideTextLine]) -> SlideTextLine | None:
    """Pick the best matching source line for one degraded Markdown line."""

    normalized_markdown = _normalize_visible_line(markdown_content).casefold()
    if not normalized_markdown:
        return None
    best_line: SlideTextLine | None = None
    best_score = 0.0
    for candidate in slide_lines:
        normalized_candidate = _normalize_visible_line(candidate.text).casefold()
        if not normalized_candidate:
            continue
        score = SequenceMatcher(None, normalized_markdown, normalized_candidate).ratio()
        if normalized_markdown and normalized_markdown[:12] in normalized_candidate:
            score += 0.08
        if len(normalized_candidate) > len(normalized_markdown):
            score += 0.04
        if score > best_score:
            best_score = score
            best_line = candidate
    if best_line is None or best_score < 0.58:
        return None
    return best_line


def _best_source_line_case_refresh(markdown_content: str, slide_lines: list[SlideTextLine]) -> SlideTextLine | None:
    """Refresh short acronym-like lines from the source when only casing is degraded."""

    normalized_markdown = _normalize_visible_line(markdown_content).casefold()
    if not normalized_markdown:
        return None
    for candidate in slide_lines:
        if normalized_markdown != _normalize_visible_line(candidate.text).casefold():
            continue
        if candidate.text == markdown_content:
            return None
        if sum(1 for ch in candidate.text if ch.isupper()) >= 2:
            return candidate
    return None


def _looks_like_degraded_pptx_line(content: str) -> bool:
    """Detect truncated or obviously damaged slide text lines."""

    if not content:
        return False
    if ELLIPSIS_RE.search(content):
        return True
    if re.search(r"[A-Za-z]\s+\.\.\.\.\s+[A-Za-z0-9]$", content):
        return True
    if len(content) >= 15 and content.count(" .") >= 2:
        return True
    return False


def _looks_like_case_softened_source_line(content: str) -> bool:
    """Detect short lines where acronym casing likely matters."""

    compact = content.strip()
    if len(compact) > 8 or " " in compact:
        return False
    return compact.isalpha() and compact[:1].isupper() and compact[1:].islower()


def _normalize_visible_line(text: str) -> str:
    """Normalize source text for matching without destroying visible meaning."""

    normalized = re.sub(r"\s+", " ", text).strip()
    normalized = normalized.replace("…", " ")
    normalized = re.sub(r"\.{3,}", " ", normalized)
    return normalized


def _normalize_markdown_line(text: str) -> str:
    """Normalize Markdown lines while stripping list and heading prefixes."""

    match = MARKDOWN_PREFIX_RE.match(text)
    content = match.group("content") if match is not None else text
    return _normalize_visible_line(content).casefold()


def _extract_pdf_text_pages(source_path: Path) -> list[str]:
    """Extract raw PDF text pages for risky OCR page recovery."""

    try:
        from pypdf import PdfReader
    except Exception:
        return []

    try:
        reader = PdfReader(str(source_path))
    except Exception:
        return []

    pages: list[str] = []
    for page in reader.pages:
        try:
            pages.append((page.extract_text() or "").strip())
        except Exception:
            pages.append("")
    return pages
